# ============================================================
# CLIMATE RISK BENCHMARKER — STREAMLIT UI (Version 6)
# Enhanced with Geography, Value Chain, Reporting Framework,
# Document Upload, and Peer Selection
# ============================================================

import streamlit as st
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_tavily import TavilySearch
from langgraph.prebuilt import create_react_agent
from datetime import datetime
import pypdf
import docx
import pptx
import io

# Load API keys
load_dotenv()

# ── Page configuration ────────────────────────────────────
st.set_page_config(
    page_title="Climate Risk Benchmarker",
    page_icon="🌍",
    layout="wide"
)

# ── Session state initialisation ──────────────────────────
# Session state remembers values across page reruns
# Without this, generated peers and manual peers
# would disappear every time any button is clicked
if "generated_peers" not in st.session_state:
    st.session_state.generated_peers = []

if "selected_generated_peers" not in st.session_state:
    st.session_state.selected_generated_peers = []

if "manual_peers" not in st.session_state:
    st.session_state.manual_peers = []

if "manual_peer_input" not in st.session_state:
    st.session_state.manual_peer_input = ""

# ── Page header ───────────────────────────────────────────
st.title("🌍 Climate Risk Benchmarker")
st.markdown("**Powered by AI** — searches peer TCFD reports, CDP disclosures, and sustainability reports to generate a climate risk long-list for any company.")
st.caption("* Required fields")
st.divider()

# ── Mandatory inputs ──────────────────────────────────────
col1, col2 = st.columns([1, 1])

with col1:
    company = st.text_input(
        "Company Name *",
        placeholder="e.g. Adelaide Airport, Tata Steel, IKEA"
    )

with col2:
    sector = st.text_input(
        "Sector / Industry *",
        placeholder="e.g. Airport Operations, Steel Manufacturing, Retail"
    )

# ── Geography ─────────────────────────────────────────────
geography = st.text_input(
    "Geography (optional)",
    placeholder="e.g. India, Australia, GCC, Global Operations"
)

# ── Value Chain Scope ─────────────────────────────────────
st.markdown("**Value Chain Scope** (optional)")
st.caption("Select which parts of the value chain to include in the analysis")

col_vc1, col_vc2, col_vc3 = st.columns([1, 1, 1])

with col_vc1:
    upstream = st.checkbox("Upstream", value=True)
with col_vc2:
    own_operations = st.checkbox("Own Operations", value=True)
with col_vc3:
    downstream = st.checkbox("Downstream", value=True)

st.divider()

# ── Reporting Framework ───────────────────────────────────
st.markdown("**Reporting Framework** (optional)")
st.caption("Select one or more frameworks to align the output to. Leave blank for a general TCFD-aligned output.")

FRAMEWORKS = {
    "TCFD": {
        "label": "TCFD",
        "description": "Task Force on Climate-related Financial Disclosures",
        "instruction": """Align the output to the TCFD framework. 
        Use TCFD's four pillars: Governance, Strategy, Risk Management, 
        and Metrics & Targets. Classify risks under TCFD categories of 
        physical risks (acute and chronic) and transition risks 
        (policy, legal, technology, market, reputational)."""
    },
    "CSRD": {
        "label": "CSRD / ESRS",
        "description": "European Sustainability Reporting Standards",
        "instruction": """Align the output to CSRD and ESRS E1 (Climate Change). 
        Apply double materiality — assess both financial materiality 
        (how climate affects the company) and impact materiality 
        (how the company affects the climate). Use ESRS E1 disclosure 
        requirements including transition plan, physical risks, 
        and climate-related targets."""
    },
    "BRSR": {
        "label": "BRSR",
        "description": "Business Responsibility and Sustainability Report (India)",
        "instruction": """Align the output to BRSR (SEBI). 
        Focus on India-specific climate risks and regulatory requirements. 
        Use BRSR's prescribed categories for environmental risks. 
        Include risks relevant to Indian regulatory context — 
        Bureau of Energy Efficiency, CPCB regulations, 
        India's NDC commitments, and PAT scheme."""
    },
    "CDP": {
        "label": "CDP",
        "description": "CDP Climate Change Questionnaire",
        "instruction": """Align the output to CDP Climate Change disclosure requirements. 
        Structure risks using CDP's risk categories and time horizons. 
        Include financial impact quantification where possible. 
        Reference CDP scoring methodology and sector-specific 
        guidance notes."""
    },
    "GRI": {
        "label": "GRI Standards",
        "description": "Global Reporting Initiative",
        "instruction": """Align the output to GRI Standards — specifically 
        GRI 201 (Economic Performance), GRI 302 (Energy), 
        GRI 303 (Water), and GRI 305 (Emissions). 
        Focus on material topics and stakeholder relevance. 
        Include both risks and opportunities framed around 
        GRI's materiality assessment approach."""
    },
    "ISSB": {
        "label": "ISSB / IFRS S2",
        "description": "International Sustainability Standards Board",
        "instruction": """Align the output to IFRS S2 Climate-related Disclosures. 
        Focus on climate-related risks and opportunities that could 
        reasonably be expected to affect the entity's cash flows, 
        access to finance, or cost of capital. Use ISSB's 
        industry-based metrics from SASB standards where relevant."""
    },
    "EU_TAXONOMY": {
        "label": "EU Taxonomy",
        "description": "EU Taxonomy for Sustainable Activities",
        "instruction": """Align the output to EU Taxonomy requirements. 
        Assess which activities substantially contribute to climate 
        change mitigation or adaptation. Identify Do No Significant Harm 
        (DNSH) criteria and minimum social safeguards. 
        Flag transition and enabling activities relevant to the sector."""
    },
    "SFDR": {
        "label": "SFDR",
        "description": "Sustainable Finance Disclosure Regulation",
        "instruction": """Align the output to SFDR Principal Adverse Impact (PAI) indicators. 
        Focus on climate and environment-related PAIs including 
        GHG emissions, carbon footprint, fossil fuel exposure, 
        and biodiversity risks. Frame risks from an 
        investor and financial product perspective."""
    }
}

fw_col1, fw_col2, fw_col3, fw_col4 = st.columns(4)
fw_columns = [fw_col1, fw_col2, fw_col3, fw_col4]

selected_frameworks = []

for i, (key, fw) in enumerate(FRAMEWORKS.items()):
    with fw_columns[i % 4]:
        if st.checkbox(
            fw["label"],
            help=fw["description"]
        ):
            selected_frameworks.append(key)

st.divider()

# ── Peer Selection ────────────────────────────────────────
st.markdown("**Peer Companies** (optional)")
st.caption("Add peers for the agent to specifically search their climate disclosures and sustainability reports.")

# ── Option 1: Generate Peers automatically ────────────────
st.markdown("**Option 1 — Generate peers automatically**")
st.caption("Click the button to let the AI suggest relevant peer companies based on your inputs above.")

generate_peers_button = st.button("🔎 Generate Peers")

# When Generate Peers is clicked
if generate_peers_button:
    if not company or not sector:
        st.warning("⚠️ Please fill in Company Name and Sector before generating peers.")
    else:
        # This is a separate mini agent call
        # It only searches for peers — not climate risks
        # Much faster than the main analysis — about 30 seconds
        with st.spinner("Finding relevant peer companies..."):

            llm_peers = ChatOpenAI(model="gpt-4o-mini", temperature=0)
            search_peers = TavilySearch(max_results=5)

            peers_agent = create_react_agent(
                model=llm_peers,
                tools=[search_peers],
                prompt="""You are a sustainability research analyst. 
                Your job is to identify relevant peer companies for 
                climate risk benchmarking. 
                Return exactly 6 peers in this format for each:
                PEER: [Company Name] | REASON: [one sentence why they are relevant]
                Only return the list. No other text."""
            )

            peers_question = f"""
            Find 6 relevant peer companies for {company}, 
            a {sector} company{f' based in {geography}' if geography else ''}.
            
            Peers should be similar in:
            - Sector and business activities
            - Geographic exposure where possible
            - Size and operational complexity
            
            Focus on companies with strong publicly available 
            climate disclosures — TCFD reports, CDP responses, 
            or sustainability reports.
            """

            peers_result = peers_agent.invoke({
                "messages": [("human", peers_question)]
            })

            peers_output = peers_result["messages"][-1].content

            # Parse the output into a list of peers
            # Each line starting with PEER: becomes one item
            parsed_peers = []
            for line in peers_output.split("\n"):
                if line.strip().startswith("PEER:"):
                    parsed_peers.append(line.strip())

            # Store in session state so they survive page reruns
            st.session_state.generated_peers = parsed_peers

# Show generated peers as checkboxes if they exist
if st.session_state.generated_peers:
    st.markdown("**Select peers to include:**")
    selected_generated_peers = []

    for peer in st.session_state.generated_peers:
        # Split the line into company name and reason
        # Format: PEER: [name] | REASON: [reason]
        try:
            parts = peer.replace("PEER:", "").split("|")
            peer_name = parts[0].strip()
            peer_reason = parts[1].replace("REASON:", "").strip() if len(parts) > 1 else ""

            # Show checkbox with peer name
            # Show reason as small grey caption below
            if st.checkbox(peer_name, key=f"gen_peer_{peer_name}"):
                selected_generated_peers.append(peer_name)
            if peer_reason:
                st.caption(f"  {peer_reason}")

        except Exception:
            # If parsing fails just show the raw line
            if st.checkbox(peer, key=f"gen_peer_{peer}"):
                selected_generated_peers.append(peer)
else:
    selected_generated_peers = []

st.divider()

# ── Option 2: Add peers manually ─────────────────────────
st.markdown("**Option 2 — Add peers manually**")
st.caption("Type any peer company name and click Add. You can add as many as you want.")

# Text input and Add button side by side
manual_col1, manual_col2 = st.columns([4, 1])

with manual_col1:
    new_peer = st.text_input(
        "Peer company name",
        placeholder="e.g. Brisbane Airport",
        label_visibility="collapsed"
    )

with manual_col2:
    add_peer_button = st.button("+ Add", type="secondary")

# When Add button is clicked
# Add the typed peer to the manual peers list in session state
if add_peer_button:
    if new_peer.strip():
        if new_peer.strip() not in st.session_state.manual_peers:
            st.session_state.manual_peers.append(new_peer.strip())
    else:
        st.warning("Please type a peer company name before clicking Add.")

# Show manually added peers with remove buttons
if st.session_state.manual_peers:
    st.markdown("**Manually added peers:**")
    for i, peer in enumerate(st.session_state.manual_peers):
        peer_col1, peer_col2 = st.columns([6, 1])
        with peer_col1:
            st.markdown(f"• {peer}")
        with peer_col2:
            # Each remove button has a unique key using the index
            if st.button("✕", key=f"remove_peer_{i}"):
                st.session_state.manual_peers.pop(i)
                st.rerun()

st.divider()

# ── Document Upload ───────────────────────────────────────
st.markdown("**Upload Documents** (optional)")
st.caption("Upload any relevant documents — previous risk assessments, sustainability reports, exposure assessments. The agent will read these alongside its web search.")
st.caption("Hold Ctrl to select multiple files at once. Accepted: PDF, PPTX, DOCX")

uploaded_files = st.file_uploader(
    "Upload files",
    accept_multiple_files=True,
    type=["pdf", "pptx", "docx"],
    label_visibility="collapsed"
)

# ── File reading functions ────────────────────────────────
def read_pdf(file):
    try:
        pdf_reader = pypdf.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Could not read PDF: {str(e)}"

def read_docx(file):
    try:
        doc = docx.Document(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        return f"Could not read DOCX: {str(e)}"

def read_pptx(file):
    try:
        presentation = pptx.Presentation(file)
        text = ""
        for slide_num, slide in enumerate(presentation.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Could not read PPTX: {str(e)}"

def extract_text_from_files(files):
    all_text = ""
    for file in files:
        file_name = file.name
        file_extension = file_name.split(".")[-1].lower()
        all_text += f"\n\n{'='*40}\n"
        all_text += f"DOCUMENT: {file_name}\n"
        all_text += f"{'='*40}\n"
        file_bytes = io.BytesIO(file.read())
        if file_extension == "pdf":
            all_text += read_pdf(file_bytes)
        elif file_extension == "docx":
            all_text += read_docx(file_bytes)
        elif file_extension == "pptx":
            all_text += read_pptx(file_bytes)
    return all_text

if uploaded_files:
    st.success(f"✅ {len(uploaded_files)} file(s) uploaded successfully:")
    for file in uploaded_files:
        st.caption(f"📄 {file.name}")

st.divider()

# ── Run button ────────────────────────────────────────────
run_button = st.button("🔍 Run Climate Risk Analysis", type="primary")

# ── Agent setup ───────────────────────────────────────────
@st.cache_resource
def setup_agent():
    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    search_broad = TavilySearch(max_results=3)
    search_deep = TavilySearch(max_results=5)
    tools = [search_broad, search_deep]

    SYSTEM_PROMPT = """You are a senior climate risk analyst with 10 years 
    of experience in TCFD-aligned risk assessments at a Big 4 consulting firm.

    Your job is to create a comprehensive long-list of climate-related risks 
    and opportunities for a company based on:
    - Peer company TCFD disclosures
    - CDP climate change responses
    - Sustainability and integrated annual reports
    - Industry association reports
    - IPCC sector-specific guidance
    - Regulatory frameworks (EU Taxonomy, CSRD, SFDR)
    - Any additional documents provided by the user

    STRICT RULES you must always follow:

    1. CATEGORIES — classify every item as one of:
       - PHYSICAL RISK — ACUTE (floods, storms, wildfires, heatwaves)
       - PHYSICAL RISK — CHRONIC (sea level rise, temperature increase, water stress)
       - TRANSITION RISK — POLICY (carbon pricing, regulations, reporting requirements)
       - TRANSITION RISK — TECHNOLOGY (low-carbon alternatives, stranded assets)
       - TRANSITION RISK — MARKET (demand shifts, commodity prices, customer behaviour)
       - TRANSITION RISK — REPUTATIONAL (brand damage, investor pressure, ESG ratings)
       - OPPORTUNITY (resource efficiency, new products, resilient supply chains)

    2. TIMEFRAME — assign one of:
       - NEAR-TERM: 0-3 years
       - MEDIUM-TERM: 3-10 years
       - LONG-TERM: 10+ years

    3. MATERIALITY — rate each as HIGH, MEDIUM, or LOW

    4. Always cite the source document or URL

    5. Format your final output exactly like this for each item:

    RISK/OPPORTUNITY #[number]
    Category: [category from list above]
    Name: [short descriptive name]
    Description: [2-3 sentences explaining the risk]
    Value Chain Affected: [Upstream / Own Operations / Downstream / Multiple]
    Timeframe: [NEAR/MEDIUM/LONG-TERM]
    Materiality: [HIGH/MEDIUM/LOW]
    Framework Relevance: [which selected frameworks flag this risk]
    Source: [document name or URL]
    ---"""

    agent = create_react_agent(
        model=llm,
        tools=tools,
        prompt=SYSTEM_PROMPT
    )

    return agent

# ── Main logic ────────────────────────────────────────────
if run_button:

    if not company or not sector:
        st.warning("⚠️ Please fill in Company Name and Sector — these are required. All other fields are optional.")

    else:
        agent = setup_agent()

        # ── Build geography context ───────────────────────
        geography_context = f"The company operates in {geography}." if geography else ""
        geography_search = f"in {geography}" if geography else ""
        geography_instruction = f"Pay particular attention to geography-specific physical risks for {geography} such as region-specific hazards, local regulations, and climate patterns." if geography else ""

        # ── Build value chain context ─────────────────────
        selected_value_chain = []
        if upstream:
            selected_value_chain.append("Upstream (raw material sourcing and supply chain)")
        if own_operations:
            selected_value_chain.append("Own Operations (direct assets, facilities, and processes)")
        if downstream:
            selected_value_chain.append("Downstream (products, customers, and end of life)")

        if selected_value_chain:
            value_chain_context = f"Focus the analysis on these value chain stages: {', '.join(selected_value_chain)}."
            value_chain_instruction = f"For each risk and opportunity, specify which value chain stage is affected — Upstream, Own Operations, or Downstream. Only include risks relevant to these selected value chain stages: {', '.join(selected_value_chain)}."
        else:
            value_chain_context = ""
            value_chain_instruction = ""

        # ── Build framework context ───────────────────────
        if selected_frameworks:
            framework_names = [FRAMEWORKS[f]["label"] for f in selected_frameworks]
            framework_instructions = "\n".join([FRAMEWORKS[f]["instruction"] for f in selected_frameworks])
            framework_context = f"""
            The output must be aligned to these reporting frameworks: {', '.join(framework_names)}.
            Framework-specific instructions:
            {framework_instructions}
            For each risk and opportunity, include a Framework Relevance field 
            showing which of these frameworks ({', '.join(framework_names)}) 
            specifically flag or require disclosure of this risk.
            """
        else:
            framework_context = "Use general TCFD-aligned categories for the output."

        # ── Build peers context ───────────────────────────
        # Combine generated peers and manual peers into one list
        all_peers = selected_generated_peers + st.session_state.manual_peers

        if all_peers:
            peers_list = ", ".join(all_peers)
            peers_context = f"""
            Specifically search the climate disclosures and sustainability 
            reports of these peer companies: {peers_list}.
            For each peer, search their:
            - TCFD disclosure or climate report
            - CDP climate change response
            - Sustainability or integrated annual report
            - Any other publicly available climate-related disclosure
            Extract relevant risks and opportunities from these specific 
            peer documents and include them in the long-list.
            """
        else:
            peers_context = ""

        # ── Extract text from uploaded documents ─────────
        document_context = ""
        if uploaded_files:
            with st.spinner("Reading uploaded documents..."):
                document_text = extract_text_from_files(uploaded_files)
                document_text_trimmed = document_text[:8000]
                document_context = f"""
                The user has also provided the following documents for analysis.
                Extract any relevant climate risks, opportunities, or exposures 
                mentioned in these documents and include them in the long-list:

                {document_text_trimmed}

                When citing risks found in these documents, reference the 
                filename as the source.
                """

        with st.spinner(f"Analysing climate risks for {company}... this takes 2-3 minutes"):

            question = f"""
            I need a comprehensive climate risk long-list for {company}, 
            a company in the {sector} sector.
            {geography_context}

            Please search for and analyse:
            1. TCFD disclosures from peer companies in the {sector} sector
            2. CDP climate change responses for {sector} companies
            3. Sustainability and annual reports from peer {sector} companies
            4. Physical climate risks specific to {sector} operations {geography_search}
            5. Transition risks from climate policy affecting {sector} {geography_search}
            6. Climate-related opportunities for {sector} companies

            Search specifically for:
            - "{sector} TCFD climate risk disclosure"
            - "{sector} CDP climate change physical risk {geography_search}"
            - "{sector} climate transition risk carbon pricing {geography_search}"
            - "{company} climate risk sustainability report"

            Create a long-list of at least 8 climate risks and opportunities.
            Follow the exact format in your instructions.
            Include both physical AND transition risks.
            {geography_instruction}
            {value_chain_context}
            {value_chain_instruction}
            {framework_context}
            {peers_context}
            {document_context}
            """

            result = agent.invoke({
                "messages": [("human", question)]
            })

            output = result["messages"][-1].content

        # ── Display results ───────────────────────────────
        st.success("✅ Analysis complete!")
        st.divider()

        st.subheader(f"Climate Risk Long-List: {company}")

        framework_display = ', '.join([FRAMEWORKS[f]["label"] for f in selected_frameworks]) if selected_frameworks else "General TCFD"
        peers_display = ', '.join(all_peers) if all_peers else "None specified"
        st.caption(f"Sector: {sector} | Geography: {geography if geography else 'Not specified'} | Frameworks: {framework_display} | Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

        if all_peers:
            st.caption(f"Peers analysed: {peers_display}")

        st.markdown(output)
        st.divider()

        # ── Download button ───────────────────────────────
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        filename = f"risk_longlist_{company.replace(' ', '_')}_{timestamp}.txt"

        download_content = f"""CLIMATE RISK LONG-LIST
Company: {company}
Sector: {sector}
Geography: {geography if geography else 'Not specified'}
Value Chain Scope: {', '.join(selected_value_chain) if selected_value_chain else 'Not specified'}
Frameworks: {framework_display}
Peers Analysed: {peers_display}
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}
{'='*60}

{output}"""

        st.download_button(
            label="📥 Download Results as Text File",
            data=download_content,
            file_name=filename,
            mime="text/plain"
        )